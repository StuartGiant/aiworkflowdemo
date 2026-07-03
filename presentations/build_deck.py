"""Build the content moderation marketing deck (3 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x23, 0x3A)   # slide background / title bar
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2D, 0x7D, 0xD2)   # blue accent
LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)   # body background
MUTED  = RGBColor(0x6B, 0x7B, 0x8D)   # secondary text
RED    = RGBColor(0xD9, 0x53, 0x4F)
GREEN  = RGBColor(0x3A, 0xA8, 0x5F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def fill_slide(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h,
    )
    shape.line.fill.background()          # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = _Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def header_bar(slide, title: str, subtitle: str = ""):
    """Dark navy bar across the top with white title."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.4), NAVY)
    add_text(slide, title,
             Inches(0.5), Inches(0.12), Inches(11), Inches(0.85),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.92), Inches(11), Inches(0.42),
                 size=14, bold=False, color=RGBColor(0xB0, 0xBE, 0xD4))


def footer(slide, label: str = "Google Chat Content Moderation  |  Cybersecurity Team"):
    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), NAVY)
    add_text(slide, label,
             Inches(0.4), Inches(7.17), Inches(10), Inches(0.3),
             size=9, color=RGBColor(0x90, 0xA0, 0xB8), align=PP_ALIGN.LEFT)


def card(slide, l, t, w, h, fill=WHITE):
    """Rounded-rect card with a very light shadow effect (two stacked rects)."""
    # shadow
    add_rect(slide, l + Inches(0.04), t + Inches(0.04), w, h,
             RGBColor(0xD0, 0xD8, 0xE8))
    return add_rect(slide, l, t, w, h, fill)


# ── Slide 1 — Problem ─────────────────────────────────────────────────────────

def slide_problem(prs: Presentation):
    slide = blank_slide(prs)
    fill_slide(slide, LIGHT)
    header_bar(slide,
               "The Problem",
               "Traditional DLP tools leave a gap in real-time communication channels")
    footer(slide)

    problems = [
        ("Delayed detection",
         "Conventional DLP operates on batch exports or end-of-day log reviews.\n"
         "By the time a violation is flagged, the data has already moved."),
        ("No evidence trail",
         "Alert-only tools produce notifications — not investigation-ready\n"
         "artefacts. There is nothing to hand to HR or Legal."),
        ("Manual triage",
         "Security analysts spend time triaging raw alerts with no context,\n"
         "no scoring, and no disposition workflow."),
    ]

    icons = ["⏱", "📂", "🔍"]   # used as simple visual markers

    col_w  = Inches(3.6)
    col_gap = Inches(0.3)
    start_l = Inches(0.55)
    top     = Inches(1.65)
    card_h  = Inches(4.7)

    for i, (title, body) in enumerate(problems):
        l = start_l + i * (col_w + col_gap)
        card(slide, l, top, col_w, card_h)

        # accent top stripe
        add_rect(slide, l, top, col_w, Inches(0.08), ACCENT)

        # icon
        add_text(slide, icons[i],
                 l + Inches(0.2), top + Inches(0.18), col_w - Inches(0.4), Inches(0.55),
                 size=28, color=ACCENT)

        # card title
        add_text(slide, title,
                 l + Inches(0.2), top + Inches(0.78), col_w - Inches(0.4), Inches(0.55),
                 size=16, bold=True, color=NAVY)

        # card body
        add_text(slide, body,
                 l + Inches(0.2), top + Inches(1.38), col_w - Inches(0.4), Inches(2.9),
                 size=13, color=MUTED, wrap=True)

    # bottom callout
    add_rect(slide, Inches(0.55), Inches(6.55), Inches(11.5), Inches(0.48),
             RGBColor(0xE8, 0xEE, 0xF8))
    add_text(slide,
             "Google Chat is a blind spot — messages are ephemeral, "
             "unmonitored in real time, and carry sensitive content.",
             Inches(0.75), Inches(6.57), Inches(11.1), Inches(0.44),
             size=12, bold=False, color=NAVY, align=PP_ALIGN.CENTER)


# ── Slide 2 — Solution ────────────────────────────────────────────────────────

def slide_solution(prs: Presentation):
    slide = blank_slide(prs)
    fill_slide(slide, LIGHT)
    header_bar(slide,
               "The Solution",
               "Real-time Google Chat content moderation — detect, act, and preserve in under 2 seconds")
    footer(slide)

    # Pipeline flow boxes
    steps = [
        ("1  Ingest",        "Google Chat\nPub/Sub listener",       NAVY),
        ("2  Detect",        "Keyword filter\n+ LLM screener\n+ Cloud Vision",  ACCENT),
        ("3  Act",           "Tombstone message\nin < 2 seconds",    RGBColor(0x21, 0x7A, 0x4E)),
        ("4  Preserve",      "Case + evidence\nin PostgreSQL",       RGBColor(0x7B, 0x3F, 0xB0)),
        ("5  Review",        "Reviewer disposition\nworkflow (DM)",  RGBColor(0xC0, 0x6C, 0x1A)),
    ]

    box_w = Inches(2.0)
    box_h = Inches(1.55)
    gap   = Inches(0.22)
    top   = Inches(1.75)
    total = len(steps) * box_w + (len(steps) - 1) * gap
    start = (SLIDE_W - total) / 2

    for i, (label, body, col) in enumerate(steps):
        l = start + i * (box_w + gap)
        add_rect(slide, l, top, box_w, box_h, col)
        add_text(slide, label,
                 l + Inches(0.1), top + Inches(0.1), box_w - Inches(0.2), Inches(0.38),
                 size=12, bold=True, color=WHITE)
        add_text(slide, body,
                 l + Inches(0.1), top + Inches(0.5), box_w - Inches(0.2), Inches(0.95),
                 size=11, color=RGBColor(0xD8, 0xE8, 0xFF), wrap=True)

        # arrow between boxes
        if i < len(steps) - 1:
            arr_l = l + box_w
            add_text(slide, "→",
                     arr_l + Inches(0.02), top + Inches(0.5),
                     gap - Inches(0.04), Inches(0.55),
                     size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # Three outcome pills
    outcomes = [
        ("BLOCK",  "Hard match → instant\nremoval, case created",   RED),
        ("REVIEW", "Borderline → human\nreviewer notified",          RGBColor(0xD9, 0x8A, 0x1A)),
        ("PASS",   "Clean content\nallowed through",                 GREEN),
    ]

    pill_w = Inches(3.1)
    pill_h = Inches(1.35)
    pill_gap = Inches(0.55)
    pill_top = Inches(3.7)
    pill_total = len(outcomes) * pill_w + (len(outcomes) - 1) * pill_gap
    pill_start = (SLIDE_W - pill_total) / 2

    add_text(slide, "Verdict routing",
             Inches(0.5), Inches(3.55), Inches(12), Inches(0.3),
             size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    for i, (label, body, col) in enumerate(outcomes):
        l = pill_start + i * (pill_w + pill_gap)
        add_rect(slide, l, pill_top, pill_w, pill_h, col)
        add_text(slide, label,
                 l + Inches(0.15), pill_top + Inches(0.1),
                 pill_w - Inches(0.3), Inches(0.42),
                 size=18, bold=True, color=WHITE)
        add_text(slide, body,
                 l + Inches(0.15), pill_top + Inches(0.52),
                 pill_w - Inches(0.3), Inches(0.75),
                 size=11, color=RGBColor(0xF8, 0xF8, 0xF8), wrap=True)

    # Key metric strip
    metrics = [
        ("< 2 sec",    "detection to action"),
        ("3-tier",     "detection layers"),
        ("100%",       "evidence preserved"),
        ("Built-in",   "human review loop"),
    ]
    strip_top = Inches(5.3)
    strip_h   = Inches(1.12)
    m_w = Inches(2.8)
    m_gap = Inches(0.2)
    m_total = len(metrics) * m_w + (len(metrics) - 1) * m_gap
    m_start = (SLIDE_W - m_total) / 2

    add_rect(slide, Inches(0), strip_top, SLIDE_W, strip_h + Inches(0.1),
             RGBColor(0xE4, 0xEA, 0xF4))

    for i, (number, label) in enumerate(metrics):
        l = m_start + i * (m_w + m_gap)
        add_text(slide, number,
                 l, strip_top + Inches(0.08), m_w, Inches(0.52),
                 size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, label,
                 l, strip_top + Inches(0.6), m_w, Inches(0.38),
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 3 — Key Capabilities ────────────────────────────────────────────────

def slide_capabilities(prs: Presentation):
    slide = blank_slide(prs)
    fill_slide(slide, LIGHT)
    header_bar(slide,
               "Key Capabilities",
               "Purpose-built for insider threat investigations")
    footer(slide)

    left_points = [
        ("Layered detection, controlled cost",
         "Hard-block keywords resolve instantly with no API cost. "
         "The LLM screener (Claude Haiku) is only invoked for ambiguous soft-flag matches. "
         "Image scoring via Cloud Vision fires only when attachments are present."),
        ("Retroactive enforcement",
         "REVIEW cases left visible pending human judgement are automatically tombstoned "
         "the moment a reviewer marks them True Positive — closing the window where "
         "content remains accessible during triage."),
        ("Investigation-ready evidence",
         "Every BLOCK and REVIEW produces a SHA-256 hashed evidence item, "
         "full scoring detail in moderation_decisions, and a timestamped case record — "
         "ready to hand to HR, Legal, or a formal IR process."),
    ]

    right_points = [
        ("Human reviewer workflow",
         "Reviewers receive a Google Chat DM with True Positive / False Positive / "
         "Inconclusive buttons. One click updates the case record. "
         "No separate tool, no email chain, no manual DB entry."),
        ("Self-healing subscription",
         "Workspace Events subscriptions expire silently. "
         "The listener auto-reactivates every 5 minutes — "
         "no operator intervention needed to maintain continuous coverage."),
        ("Auditable by design",
         "All decisions are logged with matched terms, LLM rationale, "
         "image score, final action, and engine version. "
         "Full reproducibility for post-incident review."),
    ]

    col_w  = Inches(5.6)
    left_l = Inches(0.45)
    right_l= Inches(6.85)
    top    = Inches(1.6)
    row_h  = Inches(1.6)
    row_gap= Inches(0.12)

    def point_block(l, t, title, body):
        h = row_h
        card(slide, l, t, col_w, h)
        add_rect(slide, l, t, Inches(0.06), h, ACCENT)   # left accent stripe
        add_text(slide, title,
                 l + Inches(0.2), t + Inches(0.12), col_w - Inches(0.3), Inches(0.38),
                 size=13, bold=True, color=NAVY)
        add_text(slide, body,
                 l + Inches(0.2), t + Inches(0.5), col_w - Inches(0.3), Inches(1.0),
                 size=11, color=MUTED, wrap=True)

    for i, (title, body) in enumerate(left_points):
        t = top + i * (row_h + row_gap)
        point_block(left_l, t, title, body)

    for i, (title, body) in enumerate(right_points):
        t = top + i * (row_h + row_gap)
        point_block(right_l, t, title, body)

    # divider
    add_rect(slide, Inches(6.5), Inches(1.65), Inches(0.02), Inches(4.9),
             RGBColor(0xCC, 0xD6, 0xE8))


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_problem(prs)
    slide_solution(prs)
    slide_capabilities(prs)
    out = "presentations/content_moderation_deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
